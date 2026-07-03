"""Create updated proposal Word and PowerPoint deliverables."""

from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
DELIVERABLES = ROOT / "deliverables"

TEAM_MEMBERS = [
    "Rakshita Saroha",
    "Rutvi Karnik",
    "Akanksha Madhu Kiran",
    "Dexin Yang",
    "Samyog Maharjan",
]

PROJECT_TITLE = "Dark Pattern Recognition"


def add_doc_heading(document: Document, text: str, level: int = 1) -> None:
    heading = document.add_heading(text, level=level)
    heading.paragraph_format.space_after = Pt(4)


def add_doc_bullets(document: Document, items: list[str]) -> None:
    for item in items:
        document.add_paragraph(item, style="List Bullet")


def create_docx() -> Path:
    DELIVERABLES.mkdir(exist_ok=True)
    path = DELIVERABLES / "Group 20C - Project Proposal Updated.docx"

    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    title = document.add_heading("Project Proposal", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_doc_heading(document, "Team Members & Project Title")
    document.add_paragraph(", ".join(TEAM_MEMBERS))
    document.add_paragraph(PROJECT_TITLE)

    add_doc_heading(document, "Topic & Summary")
    add_doc_heading(document, "Topic of Interest", level=2)
    document.add_paragraph(
        "Our team is trying to identify marketing dark patterns on e-commerce websites. "
        "Dark patterns are design and marketing techniques that influence users into making "
        "decisions they may not have intended, such as rushing a purchase, accepting a "
        "subscription, or sharing personal data. We plan to develop a text-based machine "
        "learning model that detects possible dark-pattern language from website copy and "
        "then use Playwright as a demo tool to collect visible text from webpages, including "
        "some popup text when it appears in the page DOM."
    )
    document.add_paragraph(
        "Our goal is to create an effective review aid that can help consumers, designers, "
        "and researchers recognize potentially manipulative e-commerce language while being "
        "clear that the model should support human review rather than make final accusations."
    )

    add_doc_heading(document, "Summary", level=2)
    add_doc_heading(document, "1. Why Is This Topic Important?", level=3)
    document.add_paragraph(
        "Dark patterns affect internet users who may not realize they are being pushed toward "
        "purchases, subscriptions, or data-sharing decisions. These tactics raise concerns "
        "about consumer privacy, digital ethics, and trust in e-commerce platforms."
    )

    add_doc_heading(document, "2. How Will Data Be Used?", level=3)
    document.add_paragraph(
        "The primary model uses a balanced labeled dataset with website text, a binary label "
        "for whether the text is a dark pattern, and a category such as Urgency, Scarcity, "
        "Social Proof, or Not Dark Pattern. We normalize the different dataset formats into "
        "common columns: text, label, category, and source. The model converts text into "
        "TF-IDF features and compares supervised learning models including Logistic "
        "Regression, Naive Bayes, Linear SVM, Decision Trees, and Random Forests."
    )
    document.add_paragraph(
        "In addition to the balanced primary dataset, we also created an expanded training "
        "experiment that incorporates Devitachi and Akash Nath data. Devitachi overlaps with "
        "the primary dataset by text, so we preserve it as source provenance rather than "
        "double-counting duplicate rows. Akash Nath contributes additional binary examples "
        "after missing text and label rows are removed."
    )

    add_doc_heading(document, "3. Potential Impact", level=3)
    document.add_paragraph(
        "The findings can help people better understand how websites use persuasive marketing "
        "language. The demo can flag suspicious phrases for closer review, helping users "
        "notice false urgency, scarcity messages, confirmshaming, and social proof claims. "
        "This could encourage more transparent design and ethical marketing practices."
    )

    add_doc_heading(document, "Machine Learning Algorithm(s)")
    add_doc_bullets(
        document,
        [
            "TF-IDF text feature extraction to convert website text into numerical features.",
            "Logistic Regression as the current primary model because it has the best F1 score and interpretable feature weights.",
            "Naive Bayes, Linear SVM, Decision Trees, and Random Forests as comparison models.",
            "Expanded dataset experiment using Devitachi and Akash Nath sources; Linear SVM performs best on this less-balanced expanded dataset.",
            "Playwright webpage scanner as a demo layer that collects visible webpage text and sends snippets to the trained text model.",
        ],
    )

    add_doc_heading(document, "Research Question")
    document.add_paragraph(
        "How effectively can text-based machine learning models identify marketing dark "
        "patterns in e-commerce website copy, and can a webpage text-scanning demo help "
        "consumers recognize manipulative marketing techniques while still preserving human judgment?"
    )

    add_doc_heading(document, "Dataset(s)")
    add_doc_bullets(
        document,
        [
            "Primary dataset: Krish Uppal, Dark Patterns, Kaggle. This balanced dataset has 2,356 examples with 1,178 dark-pattern and 1,178 non-dark-pattern text snippets.",
            "Supplemental dataset: DeVitachi, Dark-Pattern, Kaggle. This dataset includes pattern string, pattern category, pattern type, where in website, deceptive flag, and website page.",
            "Supplemental dataset: Akash Nath, Deceptive Patterns, Kaggle. Its pattern classification file contributes additional binary examples after cleaning missing rows.",
            "Expanded training option: primary + Devitachi + Akash Nath yields 4,151 unique text rows after deduplication/source aggregation.",
        ],
    )

    add_doc_heading(document, "Training & Testing Considerations")
    document.add_paragraph(
        "The primary dataset is large enough for train/test evaluation and is balanced across "
        "dark-pattern and non-dark-pattern labels, which makes binary classification more "
        "reliable than using a heavily imbalanced source alone. The project uses an 80/20 "
        "stratified split and compares models using accuracy, precision, recall, and F1 score. "
        "The current best model is TF-IDF + Logistic Regression with an F1 score around 0.936."
    )
    document.add_paragraph(
        "For the expanded dataset experiment, Linear SVM performs best with an F1 score around "
        "0.872. We keep Logistic Regression as the default app/demo model because the balanced "
        "baseline is easier to explain and Logistic Regression provides confidence scores for "
        "the Streamlit and Playwright demos."
    )

    add_doc_heading(document, "Sources of Bias")
    add_doc_bullets(
        document,
        [
            "Labeling bias: dark-pattern labels were assigned by dataset creators and may reflect subjective judgment.",
            "Category imbalance: some categories, such as Obstruction, Sneaking, and Forced Action, have far fewer examples than common categories like Scarcity and Social Proof.",
            "Domain bias: the datasets focus mostly on English-language e-commerce websites, so results may not generalize to other languages, regions, or domains.",
            "Text-only limitation: the current model does not understand page layout, colors, hidden buttons, images, or full checkout flows.",
        ],
    )

    add_doc_heading(document, "AI Impact & Ethics")
    add_doc_bullets(
        document,
        [
            "Positive impact: the tool can help consumers notice manipulative language before making purchases or sharing personal data.",
            "Positive impact: the project can help designers and companies audit website copy for potentially deceptive patterns.",
            "Risk: false positives could unfairly accuse a legitimate business or normal marketing message of being deceptive.",
            "Risk: false negatives could miss real dark patterns and leave users unprotected.",
            "Human-in-the-loop safeguard: predictions should be treated as review signals, not final judgments.",
            "Transparency safeguard: the demo reports model confidence and explains that text-only detection cannot fully evaluate visual design, checkout flow, or user intent.",
            "Future-work consideration: sentiment may help identify shame, fear, or guilt, but context such as button role, popup placement, DOM structure, and checkout step is more important for borderline cases.",
            "Production limitation: even with Playwright, the current project mainly collects text; full detection would need OCR, screenshot analysis, DOM/CSS features, and scripted user-flow testing.",
        ],
    )

    add_doc_heading(document, "Mitigation Strategies")
    add_doc_bullets(
        document,
        [
            "Clean missing or incomplete text entries before training.",
            "Use multiple datasets for exploratory analysis while keeping the primary training dataset balanced and consistent.",
            "Report precision, recall, F1 score, confusion matrix, and edge cases instead of relying only on accuracy.",
            "Use the model as a human-in-the-loop review aid, not as an automatic accusation tool.",
            "In the Playwright scanner, show only higher-confidence flags by default to reduce low-confidence false positives.",
            "Suppress low-context discount snippets and product-catalog/spec snippets during demos unless they contain urgency or scarcity pressure language.",
            "Keep sentiment and full UI-context detection as future work unless more labeled data and a more complex training pipeline are available.",
            "Evaluate OCR, color/layout analysis, and UX-flow detection separately before making production-level claims.",
        ],
    )

    add_doc_heading(document, "Citations")
    citations = [
        "Mathur, A., Acar, G., Friedman, M. J., Lucherini, E., Mayer, J., Chetty, M., & Narayanan, A. (2019). Dark Patterns at Scale: Findings from a Crawl of 11K Shopping Websites. Proceedings of the ACM on Human-Computer Interaction, 3(CSCW), 1-32. https://doi.org/10.1145/3359183",
        "Ramteke, A., Tembhurne, S., Sonawane, G., & Bhimanpallewar, R. N. (2024). Detecting Deceptive Dark Patterns in E-commerce Platforms. arXiv. https://doi.org/10.48550/arXiv.2406.01608",
        "Koh, W. C., & Seah, Y. Z. (2023). Unintended consumption: The effects of four e-commerce dark patterns. Cleaner and Responsible Consumption, 11, 100145. https://doi.org/10.1016/j.clrc.2023.100145",
        "Federal Trade Commission. (2022). Bringing Dark Patterns to Light. https://www.ftc.gov/reports/bringing-dark-patterns-light",
        "Brignull, H. (2023). Deceptive Patterns. deceptive.design. https://www.deceptive.design/",
        "Gray, C. M., Kou, Y., Battles, B., Hoggatt, J., & Toombs, A. L. (2018). The Dark (Patterns) Side of UX Design. Proceedings of the 2018 CHI Conference on Human Factors in Computing Systems, Paper 534. https://doi.org/10.1145/3173574.3174108",
        "Luguri, J., & Strahilevitz, L. J. (2021). Shining a Light on Dark Patterns. Journal of Legal Analysis, 13(1), 43-109. https://doi.org/10.1093/jla/laaa006",
    ]
    add_doc_bullets(document, citations)

    document.save(path)
    return path


def set_text_box(shape, text: str, font_size: int = 24, bold: bool = False) -> None:
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = PptPt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(33, 37, 41)


def add_slide(prs: Presentation, title: str, body: list[str] | str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = PptInches(0.6)
    title_box = slide.shapes.add_textbox(left, PptInches(0.35), PptInches(12.2), PptInches(0.7))
    set_text_box(title_box, title, font_size=30, bold=True)

    body_box = slide.shapes.add_textbox(left, PptInches(1.3), PptInches(12.0), PptInches(5.7))
    frame = body_box.text_frame
    frame.word_wrap = True
    if isinstance(body, str):
        frame.text = body
        frame.paragraphs[0].font.size = PptPt(22)
    else:
        frame.clear()
        for index, item in enumerate(body):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            paragraph.font.size = PptPt(21)
            paragraph.space_after = PptPt(8)

    footer = slide.shapes.add_textbox(PptInches(0.6), PptInches(7.0), PptInches(12.0), PptInches(0.25))
    set_text_box(footer, "Group 20C - Dark Pattern Recognition", font_size=9)


def create_pptx() -> Path:
    DELIVERABLES.mkdir(exist_ok=True)
    path = DELIVERABLES / "Group 20C - Proposal Presentation Updated.pptx"
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(PptInches(0.7), PptInches(1.6), PptInches(12), PptInches(1.2))
    set_text_box(title_box, PROJECT_TITLE, font_size=44, bold=True)
    subtitle = slide.shapes.add_textbox(PptInches(0.75), PptInches(3.0), PptInches(11.5), PptInches(1.0))
    set_text_box(subtitle, "Group Members: " + ", ".join(TEAM_MEMBERS), font_size=22)

    add_slide(
        prs,
        "Topic of Interest",
        "Our team will identify marketing dark patterns on e-commerce websites using text-based machine learning. The model detects potentially manipulative language from website copy, and our Playwright demo scans visible webpage text so users can see how the model performs on real pages.",
    )
    add_slide(
        prs,
        "1. Why Is This Topic Important?",
        "Dark patterns can pressure users into purchases, subscriptions, or data-sharing choices they did not intend. Detecting these patterns matters for consumer privacy, digital ethics, and trust in online shopping platforms.",
    )
    add_slide(
        prs,
        "2. How Will Data Be Used?",
        [
            "Normalize datasets into common columns: text, label, category, and source.",
            "Train supervised text classifiers using TF-IDF features.",
            "Compare Logistic Regression, Naive Bayes, Linear SVM, Decision Trees, and Random Forests.",
            "Run an expanded experiment with Devitachi and Akash Nath sources in addition to the balanced primary dataset.",
            "Use Playwright to collect visible webpage text during the demo and run snippets through the trained model.",
        ],
    )
    add_slide(
        prs,
        "3. Potential Impact",
        "The project can help users recognize false urgency, scarcity messages, confirmshaming, and social proof claims. It can also encourage more transparent design, while making clear that predictions should support human review rather than replace judgment.",
    )
    add_slide(
        prs,
        "Machine Learning Algorithm(s)",
        [
            "TF-IDF: converts website text into numerical features.",
            "Primary model: Logistic Regression, selected for strongest F1 score and interpretability.",
            "Comparison models: Naive Bayes, Linear SVM, Decision Tree, and Random Forest.",
            "Expanded experiment: Linear SVM performs best when adding Devitachi and Akash Nath data.",
            "Demo layer: Playwright webpage scanner for visible text and some DOM-based popup text.",
        ],
    )
    add_slide(
        prs,
        "Research Question",
        "How effectively can text-based machine learning models identify marketing dark patterns in e-commerce website copy, and can a webpage text-scanning demo help consumers recognize manipulative marketing techniques while still preserving human judgment?",
    )
    add_slide(
        prs,
        "Dataset(s)",
        [
            "Primary: Krish Uppal Dark Patterns dataset, Kaggle, 2,356 balanced examples.",
            "Supplemental: DeVitachi Dark-Pattern dataset, Kaggle, used as source provenance/metadata because text overlaps with the primary dataset.",
            "Supplemental: Akash Nath Deceptive Patterns dataset, Kaggle, adds binary classification examples after cleaning.",
            "Expanded option: 4,151 unique text rows after cleaning and deduplication/source aggregation.",
        ],
    )
    add_slide(
        prs,
        "Training & Testing Considerations",
        [
            "Use an 80/20 stratified train/test split on the balanced primary dataset.",
            "Evaluate accuracy, precision, recall, F1 score, and confusion matrix.",
            "Current best model: TF-IDF + Logistic Regression with F1 around 0.936.",
            "Expanded dataset best model: Linear SVM with F1 around 0.872.",
            "Use a confidence threshold in the webpage scanner to reduce weak false positives.",
        ],
    )
    add_slide(
        prs,
        "Sources of Bias",
        [
            "Labeling bias: labels reflect dataset creators' judgments.",
            "Category imbalance: rare categories have much less training data.",
            "Domain bias: mostly English-language e-commerce examples.",
            "Text-only limitation: layout, colors, hidden buttons, images, and checkout flows are not fully detected.",
            "Playwright can collect visible DOM text, but it is not full screenshot, color, or UX-flow understanding.",
        ],
    )
    add_slide(
        prs,
        "AI Impact & Ethics",
        [
            "Positive impact: helps consumers notice manipulative language before purchases, subscriptions, or data-sharing decisions.",
            "Positive impact: helps designers and companies audit website copy for potentially deceptive patterns.",
            "Risk: false positives could unfairly flag legitimate businesses or normal marketing language.",
            "Risk: false negatives could miss real dark patterns and leave users unprotected.",
            "Safeguard: use predictions as human-review signals, not final accusations.",
            "Transparency: the model is text-only and cannot fully judge visual layout, checkout flow, or intent.",
            "Future work: sentiment can help detect emotional pressure, but surrounding UI context is needed for reliable judgment.",
            "Production need: OCR for image text, computer vision or DOM/CSS features for color/layout, and scripted journeys for flow-based manipulation.",
        ],
    )
    add_slide(
        prs,
        "Mitigation Strategies",
        [
            "Clean missing or incomplete text entries.",
            "Use secondary datasets for comparison and exploratory analysis.",
            "Report precision, recall, F1, confusion matrix, and edge cases.",
            "Frame the model as a human-in-the-loop review aid.",
            "Filter low-confidence webpage scan results by default.",
            "Filter plain discounts and catalog-style product/spec titles unless they include pressure language.",
            "Treat sentiment, DOM structure, screenshots, and full UI flow as future work.",
            "Separate text, OCR, visual, and flow evaluation so each layer has clear evidence.",
        ],
    )
    add_slide(
        prs,
        "Citations",
        [
            "Mathur et al. (2019). Dark Patterns at Scale. https://doi.org/10.1145/3359183",
            "Ramteke et al. (2024). Detecting Deceptive Dark Patterns in E-commerce Platforms. https://doi.org/10.48550/arXiv.2406.01608",
            "Koh & Seah (2023). Unintended consumption. https://doi.org/10.1016/j.clrc.2023.100145",
            "FTC (2022). Bringing Dark Patterns to Light.",
            "Brignull (2023). Deceptive Patterns. https://www.deceptive.design/",
            "Gray et al. (2018). The Dark (Patterns) Side of UX Design. https://doi.org/10.1145/3173574.3174108",
            "Luguri & Strahilevitz (2021). Shining a Light on Dark Patterns. https://doi.org/10.1093/jla/laaa006",
        ],
    )

    thanks = prs.slides.add_slide(prs.slide_layouts[6])
    thanks_box = thanks.shapes.add_textbox(PptInches(0.7), PptInches(2.7), PptInches(12), PptInches(1.0))
    set_text_box(thanks_box, "Thank you!", font_size=44, bold=True)
    thanks_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(path)
    return path


def main() -> None:
    docx_path = create_docx()
    pptx_path = create_pptx()
    print(f"Created {docx_path}")
    print(f"Created {pptx_path}")


if __name__ == "__main__":
    main()
