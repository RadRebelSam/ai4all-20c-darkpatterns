"""Update final presentation and poster with the two-model project results."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FINAL_PRESENTATION = ROOT / "deliverables" / "Group 20C - Final Presentation.pptx"
FINAL_POSTER = ROOT / "deliverables" / "Group 20C - Final Poster.pptx"
FIGURE_DIR = ROOT / "reports" / "figures"

AUTO_SLIDE_TITLES = {
    "Current Two-Model System",
    "Second-Stage Category Results",
    "Category Error Analysis",
    "Example Result Flow",
}
AUTO_POSTER_TEXT = {
    "Two-Model Demo Flow",
    "Updated: Two Trained Models",
    "Model 1 detects Dark vs Not Dark (F1 0.936). Model 2 predicts likely type with Linear SVM (macro F1 0.894).",
}


def replace_text_in_shape(shape, replacements: dict[str, str]) -> None:
    if not hasattr(shape, "text_frame"):
        return
    if hasattr(shape, "text"):
        full_text = shape.text
        replaced_text = full_text
        for old, new in replacements.items():
            replaced_text = replaced_text.replace(old, new)
        if replaced_text != full_text:
            shape.text = replaced_text
            return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text
            for old, new in replacements.items():
                text = text.replace(old, new)
            run.text = text


def replace_text_everywhere(prs: Presentation, replacements: dict[str, str]) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)


def slide_title(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(" ".join(shape.text.split()))
    return texts[0] if texts else ""


def remove_auto_slides(prs: Presentation) -> None:
    slide_ids_to_remove = []
    for slide in prs.slides:
        title = slide_title(slide)
        if title in AUTO_SLIDE_TITLES:
            slide_ids_to_remove.append(slide.slide_id)

    if not slide_ids_to_remove:
        return

    slide_id_list = prs.slides._sldIdLst  # python-pptx has no public delete API.
    for slide_id in slide_ids_to_remove:
        for rel_id, slide in list(prs.part._rels.items()):
            target = slide._target
            if (
                "slide" in slide.reltype
                and hasattr(target, "slide_id")
                and target.slide_id == slide_id
            ):
                prs.part.drop_rel(rel_id)
                break
        for sld_id in list(slide_id_list):
            if int(sld_id.id) == slide_id:
                slide_id_list.remove(sld_id)
                break


def set_text(shape, text: str, size: int = 24, bold: bool = False, color=(33, 37, 41)) -> None:
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(*color)


def add_title(slide, title: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
    set_text(title_box, title, size=27, bold=True, color=(15, 23, 42))


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.size = Pt(size)
        paragraph.space_after = Pt(7)


def add_footer(slide) -> None:
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12), Inches(0.25))
    set_text(footer, "Group 20C - Dark Pattern Recognition", size=9, color=(100, 116, 139))


def add_picture(slide, filename: str, left: float, top: float, width: float | None = None, height: float | None = None) -> None:
    path = FIGURE_DIR / filename
    if not path.exists():
        raise FileNotFoundError(path)
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def remove_auto_poster_overlays(slide) -> None:
    for shape in list(slide.shapes):
        if hasattr(shape, "text") and shape.text.strip() in AUTO_POSTER_TEXT:
            shape._element.getparent().remove(shape._element)


def add_model_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Current Two-Model System")
    add_picture(slide, "two_stage_pipeline.png", 0.65, 1.15, width=11.9)
    add_bullets(
        slide,
        [
            "Model 1 answers whether the text looks suspicious at all.",
            "Model 2 runs only after a suspicious result and predicts a likely dark-pattern type.",
            "This keeps the demo understandable: first detection, then explanation.",
        ],
        0.9,
        5.65,
        11.4,
        1.15,
        size=16,
    )
    add_footer(slide)


def add_category_results_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Second-Stage Category Results")
    add_picture(slide, "category_model_comparison.png", 0.55, 1.05, width=5.9)
    add_picture(slide, "category_per_class_f1.png", 6.85, 1.05, width=5.8)
    add_bullets(
        slide,
        [
            "Best category model: TF-IDF + Linear SVM.",
            "Macro F1 is about 0.894, which is useful for explanation but still limited by category imbalance.",
            "Rare categories such as Sneaking, Obstruction, and Forced Action should be treated carefully.",
        ],
        0.75,
        6.25,
        12.0,
        0.75,
        size=14,
    )
    add_footer(slide)


def add_category_error_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Category Error Analysis")
    add_picture(slide, "category_confusion_matrix.png", 0.65, 1.0, width=6.1)
    add_bullets(
        slide,
        [
            "The second model is strongest on categories with more training examples.",
            "Some categories overlap semantically, especially pressure-based patterns such as urgency and scarcity.",
            "For the presentation, describe the output as a likely type, not a final judgment.",
        ],
        7.15,
        1.45,
        5.3,
        3.0,
        size=19,
    )
    add_picture(slide, "two_model_summary_card.png", 6.95, 4.75, width=5.45)
    add_footer(slide)


def add_example_flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Example Result Flow")
    add_picture(slide, "example_result_flow.png", 0.85, 1.1, width=11.4)
    add_bullets(
        slide,
        [
            "This is how the Streamlit app and scanner should explain a result: main answer first, likely type second.",
            "The category result adds interpretability without pretending to understand the full UI context.",
        ],
        1.0,
        5.95,
        11.2,
        0.85,
        size=16,
    )
    add_footer(slide)


def update_final_presentation() -> Path:
    prs = Presentation(FINAL_PRESENTATION)
    remove_auto_slides(prs)
    replace_text_everywhere(
        prs,
        {
            "MOVED from broad exploration to a balanced primary baseline plus an expanded robustness experiment across multiple classifiers.": "Moved from broad exploration to a balanced binary baseline plus a second-stage category classifier.",
            "A classic yet powerful NLP pipeline: TF-IDF vectorization paired with Logistic Regression for supervised binary text classification.": "A two-stage NLP pipeline: TF-IDF + Logistic Regression for binary detection, then TF-IDF + Linear SVM for likely type recognition.",
            "Multi-Class Classification": "Expanded Category Recognition",
            "Identify the type of dark pattern (Urgency, Scarcity, Misdirection…) rather than just dark vs. not - enabling more targeted consumer guidance.": "Improve the current second-stage category model with more balanced category labels and stronger UI context.",
            "Model Evaluation": "Model 1 Evaluation",
            "Data Visualizations": "Visual Evidence",
        },
    )
    add_model_slide(prs)
    add_category_results_slide(prs)
    add_category_error_slide(prs)
    add_example_flow_slide(prs)
    prs.save(FINAL_PRESENTATION)
    return FINAL_PRESENTATION


def update_final_poster() -> Path:
    prs = Presentation(FINAL_POSTER)
    slide = prs.slides[0]
    remove_auto_poster_overlays(slide)
    replace_text_everywhere(
        prs,
        {
            "Model\nTF-IDF + Logistic Regression. Compared against Linear SVM, Naive Bayes, Decision Tree, and Random Forest.": "Models\nModel 1: TF-IDF + Logistic Regression detects Dark vs Not Dark. Model 2: TF-IDF + Linear SVM predicts likely type.",
            "Result\nBest primary model: 93.9% accuracy, 0.936 F1.": "Results\nModel 1: 93.9% accuracy, 0.936 F1. Model 2: 0.894 macro F1 for category/type recognition.",
            "Model TF-IDF + Logistic Regression. Compared against Linear SVM, Naive Bayes, Decision Tree, and Random Forest.": "Models Model 1: TF-IDF + Logistic Regression detects Dark Pattern vs Not Dark Pattern. Model 2: TF-IDF + Linear SVM predicts likely type.",
            "Result Best primary model: 93.9% accuracy, 0.936 F1.": "Results Model 1: 93.9% accuracy, 0.936 F1. Model 2: 0.894 macro F1 for category/type recognition.",
            "Future Work Hard-negative mining, TF-IDF tuning, threshold tuning, OCR, visual layout features, and UX-flow": "Future Work Hard-negative mining, more balanced category labels, threshold tuning, OCR, visual layout features, and UX-flow",
        },
    )

    # Add compact visual panels. These supplement the existing poster even when
    # the original text is split across PowerPoint runs and cannot be replaced
    # reliably with exact string matching.
    headline = slide.shapes.add_textbox(Inches(5.95), Inches(0.45), Inches(2.85), Inches(0.35))
    set_text(headline, "Updated: Two Trained Models", size=12, bold=True, color=(15, 23, 42))

    summary = slide.shapes.add_textbox(Inches(5.95), Inches(0.82), Inches(2.85), Inches(0.82))
    set_text(
        summary,
        "Model 1 detects Dark vs Not Dark (F1 0.936). Model 2 predicts likely type with Linear SVM (macro F1 0.894).",
        size=8,
        color=(30, 41, 59),
    )

    panel = slide.shapes.add_textbox(Inches(6.75), Inches(4.15), Inches(2.35), Inches(0.34))
    set_text(panel, "Two-Model Demo Flow", size=12, bold=True, color=(15, 23, 42))
    add_picture(slide, "two_model_summary_card.png", 6.45, 4.48, width=2.65)

    prs.save(FINAL_POSTER)
    return FINAL_POSTER


def main() -> None:
    presentation = update_final_presentation()
    poster = update_final_poster()
    print(f"Updated {presentation}")
    print(f"Updated {poster}")


if __name__ == "__main__":
    main()
